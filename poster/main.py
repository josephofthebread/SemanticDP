import math
import string
from typing import TypedDict

import matplotlib
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import pyphen
import qrcode
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.slide import Slide
from pptx.text.text import TextFrame
from pptx.util import Cm, Emu, Pt

matplotlib.use("Agg")


class Block(TypedDict, total=False):
  kind: str
  text: str
  bold: bool
  bullet: bool


INK_HEX = "#1a1a2e"
GRID_HEX = "#d9dce3"
MECH: dict[str, str] = {
  "Non-private": "#8a8f9a",
  "Random replacement": "#e0662b",
  "Sanitization": "#2f7fd1",
  "Private optimization": "#3a9e6f",
}
ACCENT = RGBColor(0x2F, 0x7F, 0xD1)
INK = RGBColor(0x1A, 0x1A, 0x2E)
GREY = RGBColor(0x55, 0x59, 0x63)


def style() -> None:
  plt.rcParams.update(
    {
      "figure.dpi": 200,
      "savefig.dpi": 200,
      "font.size": 15,
      "font.family": "DejaVu Sans",
      "axes.edgecolor": INK_HEX,
      "axes.labelcolor": INK_HEX,
      "axes.titlecolor": INK_HEX,
      "xtick.color": INK_HEX,
      "ytick.color": INK_HEX,
      "text.color": INK_HEX,
      "axes.linewidth": 1.1,
      "axes.grid": True,
      "grid.color": GRID_HEX,
      "grid.linewidth": 0.9,
      "axes.spines.top": False,
      "axes.spines.right": False,
      "figure.facecolor": "white",
      "axes.facecolor": "white",
      "legend.frameon": False,
    }
  )


def fig_coefficient() -> None:
  comps = [
    ("Entity\nfidelity", 0.733, 0.122),
    ("Instruction\nfollowing", 0.524, 0.090),
    ("Composite\nalignment", 0.436, 0.077),
  ]
  est = np.array([c[1] for c in comps])
  se = np.array([c[2] for c in comps])
  y = np.arange(len(comps))[::-1]
  fig, ax = plt.subplots(figsize=(8.8, 4.4))
  ax.axvspan(0, 0.10, color="#eeeef1", zorder=0)
  ax.axvline(0.10, color=INK_HEX, lw=1.6, ls=(0, (4, 3)), zorder=2)
  ax.errorbar(
    est, y, xerr=1.96 * se, fmt="o", ms=14, lw=3.2, color="#2f7fd1", ecolor="#2f7fd1", capsize=7, capthick=3, zorder=3
  )
  for xi, yi in zip(est, y):
    ax.annotate(
      f"{xi:.2f}", (xi, yi), textcoords="offset points", xytext=(0, 17), ha="center", fontsize=16, fontweight="bold"
    )
  ax.set_yticks(y)
  ax.set_yticklabels([c[0] for c in comps], fontsize=15.5)
  ax.set_ylim(-0.7, len(comps) - 0.25)
  ax.set_xlim(0, 0.95)
  ax.set_xlabel("Standardized effect of semantic distortion,  |c|", fontsize=15)
  ax.text(0.05, 1.0, "below\nmeaningful\neffect", rotation=90, ha="center", va="center", fontsize=11.5, color="#8a8f9a")
  ax.annotate(
    "meaningful-effect bar",
    xy=(0.10, len(comps) - 0.42),
    xytext=(0.30, len(comps) - 0.30),
    fontsize=12,
    color=INK_HEX,
    ha="left",
    va="center",
    arrowprops={"arrowstyle": "-", "color": INK_HEX, "lw": 1},
  )
  ax.set_title("Distortion drives every capability", fontsize=17, fontweight="bold", pad=14)
  ax.grid(axis="y", visible=False)
  fig.subplots_adjust(top=0.85, bottom=0.20, left=0.16, right=0.97)
  fig.savefig("fig/fig1_coefficient.png")
  plt.close(fig)


def fig_cliff() -> None:
  eps = [1, 3, 6, 12]
  distortion = [0.4598, 0.4572, 0.0833, 0.0]
  fidelity = [0.705, 0.700, 0.950, 0.975]
  fig, ax1 = plt.subplots(figsize=(8.6, 4.9))
  x = np.arange(4)
  ax1.bar(x, distortion, width=0.62, color="#2f7fd1", alpha=0.28, edgecolor="#2f7fd1", lw=1.5, zorder=1)
  ax1.set_ylabel("Semantic distortion $S$", color="#2f7fd1", fontsize=15)
  ax1.tick_params(axis="y", labelcolor="#2f7fd1")
  ax1.set_ylim(0, 0.72)
  ax1.set_xticks(x)
  ax1.set_xticklabels([f"ε={e}" for e in eps], fontsize=15)
  ax1.set_xlabel("Privacy budget  (strict → loose)", fontsize=15)
  ax1.grid(False)
  ax2 = ax1.twinx()
  ax2.plot(x, fidelity, "-o", color="#3a9e6f", lw=3.4, ms=12, zorder=3)
  ax2.set_ylabel("Entity fidelity", color="#3a9e6f", fontsize=15)
  ax2.tick_params(axis="y", labelcolor="#3a9e6f")
  ax2.set_ylim(0.6, 1.02)
  ax2.grid(False)
  ax2.spines["top"].set_visible(False)
  ax1.plot([-0.30, 1.30], [0.60, 0.60], color=INK_HEX, lw=1, ls=":")
  ax1.text(0.5, 0.66, "budget triples,\nnothing moves", fontsize=12, ha="center", va="center", color=INK_HEX)
  ax1.annotate("", xy=(1.5, 0.10), xytext=(1.5, 0.50), arrowprops={"arrowstyle": "-|>", "color": INK_HEX, "lw": 2.2})
  ax1.text(1.63, 0.32, "the cliff:\n3× tighter budget,\nsame distortion", fontsize=12, va="center", ha="left")
  ax1.set_title("Outcomes follow distortion, not the budget", fontsize=16.5, fontweight="bold", pad=12)
  fig.tight_layout()
  fig.savefig("fig/fig2_cliff.png")
  plt.close(fig)


def fig_pathways() -> None:
  ev = pd.read_csv("../results/evaluate.csv")
  name = {"m0": "Non-private", "m1": "Random replacement", "m2": "Sanitization", "m3": "Private optimization"}
  ev["mech"] = ev["mechanism"].map(name)
  nem = ev[(ev.corpus == "nemotron") & (ev.mechanism.isin(["m1", "m2", "m3"]))]
  fig, axes = plt.subplots(1, 2, figsize=(12.6, 4.9))
  metrics = [("entity_fidelity/correct", "Entity fidelity"), ("ifeval/prompt_accuracy", "Instruction following")]
  for ax, (metric, title) in zip(axes, metrics):
    base = ev[(ev.mechanism == "m0") & (ev.corpus == "nemotron")][metric].mean()
    ax.axhline(base, color=MECH["Non-private"], lw=2.2, ls=(0, (5, 3)), zorder=1)
    ax.text(4.02, base, "non-private", color=MECH["Non-private"], fontsize=11, va="center", ha="left")
    for mech in ["Random replacement", "Sanitization", "Private optimization"]:
      series = nem[nem.mech == mech].groupby("d")[metric].mean()
      ax.plot(series.index, series.values, "-o", color=MECH[mech], lw=3, ms=10, label=mech, zorder=3)
    ax.set_title(title, fontsize=16, fontweight="bold")
    ax.set_xlabel("Privacy-level rank  (1 = strictest)", fontsize=14)
    ax.set_xticks([1, 2, 3, 4])
    ax.set_xlim(0.8, 4.9)
  axes[0].set_ylabel("score", fontsize=14)
  handles, labels = axes[0].get_legend_handles_labels()
  fig.legend(handles, labels, loc="lower center", ncol=3, fontsize=13.5, bbox_to_anchor=(0.5, -0.06))
  fig.suptitle("Two mechanisms, two distinct failures", fontsize=17, fontweight="bold", y=1.02)
  fig.tight_layout()
  fig.savefig("fig/fig3_pathways.png", bbox_inches="tight")
  plt.close(fig)


def fig_fabrication() -> None:
  fab = {
    "Random replacement": [0.0643, 0.0184, 0.0100, 0.0060],
    "Sanitization": [0.1727, 0.2011, 0.0214, 0.0033],
    "Private optimization": [0.0377, 0.0214, 0.0163, 0.0124],
  }
  baseline = 0.0042
  fig, ax = plt.subplots(figsize=(8.8, 4.9))
  x = np.arange(4)
  width = 0.26
  for i, (mech, vals) in enumerate(fab.items()):
    ax.bar(x + (i - 1) * width, vals, width, color=MECH[mech], label=mech, edgecolor="white", lw=0.8)
  ax.axhline(baseline, color=INK_HEX, lw=1.6, ls=(0, (4, 3)))
  ax.annotate(
    "non-private baseline  0.4%",
    xy=(3.15, baseline),
    xytext=(3.4, 0.055),
    fontsize=11.5,
    ha="right",
    va="bottom",
    color=INK_HEX,
    arrowprops={"arrowstyle": "-", "color": INK_HEX, "lw": 0.8},
  )
  ax.annotate(
    "48× baseline",
    xy=(1, 0.2011),
    xytext=(1.15, 0.215),
    fontsize=13,
    fontweight="bold",
    color="#2f7fd1",
    arrowprops={"arrowstyle": "-|>", "color": "#2f7fd1", "lw": 2},
  )
  ax.set_xticks(x)
  ax.set_xticklabels(["rank 1\n(strict)", "rank 2", "rank 3", "rank 4\n(loose)"], fontsize=13)
  ax.set_ylabel("Fabrication rate", fontsize=15)
  ax.set_ylim(0, 0.235)
  ax.legend(fontsize=12.5, loc="upper right")
  ax.set_title("Fabrication is driven by text rewriting, not gradient noise", fontsize=15.5, fontweight="bold", pad=10)
  fig.tight_layout()
  fig.savefig("fig/fig4_fabrication.png")
  plt.close(fig)


def make_equations() -> None:
  plt.rcParams["mathtext.fontset"] = "cm"
  forms = {
    "eq_reg": r"$Y_i \;=\; a_m \,+\, b_m\, d_i \,+\, c\, S_i \,+\, e_i$",
    "eq_sdist": r"$S \;=\; \mathbb{E}\,[\,1 - \cos(\phi(x),\, \phi(\tilde{x}))\,]$",
    "eq_tem": r"$\mathbb{P}(x\to w)\ \propto\ \exp\!\left(\frac{\varepsilon}{2}\,u(w)\right)$",
  }
  for name, tex in forms.items():
    fig = plt.figure(figsize=(9, 1.4))
    fig.text(0.5, 0.5, tex, fontsize=34, color=INK_HEX, ha="center", va="center")
    fig.savefig(f"fig/{name}.png", dpi=300, transparent=True, bbox_inches="tight", pad_inches=0.06)
    plt.close(fig)


def make_qr() -> None:
  qr = qrcode.QRCode(version=1, box_size=20, border=2, error_correction=qrcode.constants.ERROR_CORRECT_M)
  qr.add_data("https://github.com/josephofthebread/SemanticDP")
  qr.make(fit=True)
  qr.make_image(fill_color=INK_HEX, back_color="white").save("fig/qr.png")


def make_figures() -> None:
  style()
  fig_coefficient()
  fig_cliff()
  fig_pathways()
  fig_fabrication()
  make_equations()
  make_qr()


DICT = pyphen.Pyphen(lang="en_US")
SHY = "­"
NBSP = " "
LX = 1.8
RX = 31.1
CW = 27.0
BODY_TOP = 12.8
GAP = 0.2
HEAD_GAP = 0.4
LINE_SPACING = 0.9
SPACE_AFTER = 2
SZ: dict[str, int] = {"head": 31, "sub": 24, "body": 20, "ref": 14, "cap": 15}


def hyphenate_word(word: str) -> str:
  if "-" in word:
    return "-".join(hyphenate_word(p) for p in word.split("-"))
  core = word.strip(string.punctuation + "“”’‘—…")
  if len(core) >= 7 and core.isalpha() and not core.isupper():
    return word.replace(core, DICT.inserted(core, SHY), 1)
  return word


def hyphenate(text: str) -> str:
  return " ".join(hyphenate_word(w) for w in text.split(" "))


def chars_per_line(size: int) -> int:
  return max(14, int(CW / (size * 0.0182)))


def line_cm(size: int) -> float:
  return size / 72 * 2.54 * 1.2 * LINE_SPACING


def estimate_lines(text: str, size: int) -> int:
  return max(1, math.ceil(len(text) / chars_per_line(size)))


def panel_height(blocks: list[Block]) -> float:
  height = 0.12
  for block in blocks:
    size = SZ[block["kind"]]
    height += estimate_lines(block["text"], size) * line_cm(size) + SPACE_AFTER / 72 * 2.54
  return height + 0.10


def fill_frame(frame: TextFrame, blocks: list[Block]) -> None:
  frame.word_wrap = True
  frame.auto_size = MSO_AUTO_SIZE.NONE
  frame.margin_left = Cm(0.06)
  frame.margin_right = Cm(0.06)
  frame.margin_top = Cm(0.03)
  frame.margin_bottom = Cm(0.03)
  for i, block in enumerate(blocks):
    paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
    kind = block["kind"]
    justify = kind not in ("head", "sub", "ref") and not block.get("bullet")
    paragraph.alignment = PP_ALIGN.JUSTIFY if justify else PP_ALIGN.LEFT
    paragraph.line_spacing = LINE_SPACING
    paragraph.space_after = Pt(SPACE_AFTER)
    paragraph.space_before = Pt(0)
    run = paragraph.add_run()
    text = block["text"]
    if kind in ("body", "ref", "cap"):
      text = hyphenate(text)
    if kind == "body" and " " in text:
      head, _, tail = text.rpartition(" ")
      text = head + NBSP + tail
    run.text = ("•  " + text) if block.get("bullet") else text
    font = run.font
    font.size = Pt(SZ[kind])
    font.bold = block.get("bold", kind in ("head", "sub"))
    font.color.rgb = ACCENT if kind == "sub" else INK


def add_panel(slide: Slide, x: float, y: float, blocks: list[Block]) -> float:
  height = panel_height(blocks)
  box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(CW), Cm(height))
  fill_frame(box.text_frame, blocks)
  return y + height


def add_figure(slide: Slide, x: float, y: float, path: str, w: float) -> float:
  iw, ih = Image.open(path).size
  height = w * ih / iw
  slide.shapes.add_picture(path, Cm(x + (CW - w) / 2), Cm(y), Cm(w), Cm(height))
  return y + height


def add_caption(slide: Slide, x: float, y: float, text: str) -> float:
  text = hyphenate(text)
  height = estimate_lines(text, SZ["cap"]) * line_cm(SZ["cap"]) + 0.15
  box = slide.shapes.add_textbox(Cm(x), Cm(y + 0.05), Cm(CW), Cm(height))
  frame = box.text_frame
  frame.word_wrap = True
  frame.auto_size = MSO_AUTO_SIZE.NONE
  frame.margin_top = Cm(0.02)
  frame.margin_bottom = Cm(0.02)
  paragraph = frame.paragraphs[0]
  paragraph.line_spacing = LINE_SPACING
  run = paragraph.add_run()
  run.text = text
  run.font.size = Pt(SZ["cap"])
  run.font.italic = True
  run.font.color.rgb = INK
  return y + height + 0.05


def head(text: str) -> Block:
  return {"kind": "head", "text": text}


def sub(text: str) -> Block:
  return {"kind": "sub", "text": text}


def body(text: str, bold: bool = False, bullet: bool = False) -> Block:
  return {"kind": "body", "text": text, "bold": bold, "bullet": bullet}


def ref(text: str) -> Block:
  return {"kind": "ref", "text": text}


def set_title(frame: TextFrame) -> None:
  for paragraph in list(frame.paragraphs):
    paragraph._p.getparent().remove(paragraph._p)
  first = frame.add_paragraph()
  run = first.add_run()
  run.text = "Does Privacy Distort Alignment?"
  run.font.size = Pt(54)
  run.font.bold = True
  run.font.color.rgb = INK
  second = frame.add_paragraph()
  run = second.add_run()
  run.text = "Why privacy mechanisms that rewrite text make LLMs hallucinate — and how to mitigate it"
  run.font.size = Pt(24)
  run.font.color.rgb = ACCENT


def add_author(slide: Slide, x: float, name: str, affiliation: str) -> None:
  box = slide.shapes.add_textbox(Cm(x), Cm(8.7), Cm(15), Cm(2.4))
  frame = box.text_frame
  frame.word_wrap = True
  frame.auto_size = MSO_AUTO_SIZE.NONE
  name_paragraph = frame.paragraphs[0]
  name_paragraph.alignment = PP_ALIGN.CENTER
  run = name_paragraph.add_run()
  run.text = name
  run.font.size = Pt(30)
  run.font.bold = True
  run.font.color.rgb = INK
  affiliation_paragraph = frame.add_paragraph()
  affiliation_paragraph.alignment = PP_ALIGN.CENTER
  affiliation_paragraph.line_spacing = 0.95
  run = affiliation_paragraph.add_run()
  run.text = affiliation
  run.font.size = Pt(22)
  run.font.color.rgb = GREY


def left_column(slide: Slide) -> None:
  y = BODY_TOP
  y = (
    add_panel(
      slide,
      LX,
      y,
      [
        head("Motivation"),
        body(
          "Language models are increasingly fine-tuned on sensitive data — clinical notes, financial records. "
          "Differential privacy (DP) is the standard guarantee against its disclosure."
        ),
        body(
          "Yet DP fine-tuning increases hallucination, and prior work established the effect without identifying its "
          "cause — attributing it only loosely to gradient clipping and added noise."
        ),
        body("Two candidate causes that a single privacy parameter cannot separate:"),
        body("DP corrupts the meaning of the training text; or", bullet=True),
        body("it perturbs only the gradients, leaving the text intact.", bullet=True),
        body("Which cause dominates determines whether improved mechanism design can recover alignment."),
      ],
    )
    + HEAD_GAP
  )
  y = (
    add_panel(
      slide,
      LX,
      y,
      [head("Problem"), body("We distinguish the two causes with a single pre-registered regression:")],
    )
    + 0.14
  )
  y = add_figure(slide, LX, y, "fig/eq_reg.png", 10) + 0.16
  y = (
    add_panel(
      slide,
      LX,
      y,
      [
        body(
          "Y is alignment degradation and d the privacy level; S is semantic distortion — the extent to which a "
          "mechanism alters the meaning of the text, computed without reference to the budget ε:"
        )
      ],
    )
    + 0.10
  )
  y = add_figure(slide, LX, y, "fig/eq_sdist.png", 10) + 0.16
  y = (
    add_panel(
      slide,
      LX,
      y,
      [
        body(
          "A large, dominant coefficient c — its confidence interval excluding the pre-registered minimum effect — "
          "identifies meaning as the operative cause."
        ),
        body(
          "Identification: S is varied at matched privacy by contrasting mechanism families, together with a "
          "non-private control carrying no guarantee — so that DP-specific effects are tested rather than assumed."
        ),
      ],
    )
    + HEAD_GAP
  )
  y = (
    add_panel(
      slide,
      LX,
      y,
      [
        head("Method"),
        body("A single fixed LoRA (low-rank) fine-tuning procedure; only the privacy mechanism varies:"),
        body("none — the non-private baseline.", bullet=True),
        body("random word replacement — semantic damage with no privacy guarantee (the control).", bullet=True),
        body("meaning-preserving sanitization — word-level metric-DP over GloVe embeddings (TEM).", bullet=True),
        body("DP-SGD — calibrated Gaussian noise on the gradients, ε ∈ {1, 4, 8, 16}.", bullet=True),
        body(
          "Holding the procedure fixed makes the mechanism the only variable. All metrics are judge-free — no LLM "
          "grader whose own alignment could contaminate the measurement."
        ),
      ],
    )
    + HEAD_GAP
  )
  y = (
    add_panel(
      slide,
      LX,
      y,
      [
        head("Dataset"),
        body("Two datasets differing in entity density, where privacy risk concentrates:"),
        body(
          "Nemotron-PII — synthetic clinical records with gold entity spans, so hallucination and entity distortion "
          "are scored precisely against ground truth.",
          bullet=True,
        ),
        body("Alpaca — generic instruction tasks, few named entities.", bullet=True),
        body(
          "3 model families (Qwen3-1.7B, Llama-3.2-1B, Falcon3-1B) × 3 seeds × 13 privacy settings × 2 datasets = "
          "234 adapters."
        ),
      ],
    )
    + HEAD_GAP
  )
  y = (
    add_panel(
      slide,
      LX,
      y,
      [
        sub("Two mechanisms, two distinct failures"),
        body(
          "Text corruption degrades instruction-following; DP-SGD degrades factual recall without altering the text "
          "at all — its text distortion is exactly zero. A single degradation axis cannot describe both."
        ),
      ],
    )
    + GAP
  )
  y = add_figure(slide, LX, y, "fig/fig3_pathways.png", 17) + 0.02
  y = (
    add_caption(
      slide,
      LX,
      y,
      "DP-SGD reduces factual recall with no change to the text; the text mechanisms reduce instruction-following",
    )
    + HEAD_GAP
  )
  y = (
    add_panel(
      slide,
      LX,
      y,
      [
        head("Distorted meaning → fabricated facts"),
        body(
          "When meaning is preserved the model remains grounded; when it is destroyed the model fabricates. "
          "Metric-LDP sanitization produces content absent from the record at up to 48× the non-private baseline "
          "(ε=3); random replacement and DP-SGD fabricate 3–5× less (15× and 9× baseline), at comparable overall "
          "accuracy."
        ),
      ],
    )
    + GAP
  )
  y = add_figure(slide, LX, y, "fig/fig4_fabrication.png", 16) + 0.02
  add_caption(slide, LX, y, "Fabrication rises precisely where the text’s meaning is destroyed")


def right_column(slide: Slide) -> None:
  y = BODY_TOP
  y = (
    add_panel(
      slide,
      RX,
      y,
      [
        head("Key finding"),
        body("Privacy degrades alignment through meaning, not noise.", bold=True),
        body(
          "Semantic distortion is the dominant driver of degradation — five times the pre-registered minimum effect "
          "(p < 10⁻⁴), and larger than the privacy budget for every mechanism. It is identified from the pooled "
          "design across mechanism families, not from any single one."
        ),
        body(
          "A third pre-registered hypothesis — moderation by entity density — is not supported: 95% CI "
          "[−0.03, +0.24] includes zero.",
          bullet=True,
        ),
      ],
    )
    + GAP
  )
  y = add_figure(slide, RX, y, "fig/fig1_coefficient.png", 21) + 0.02
  y = (
    add_caption(slide, RX, y, "Every capability lies well beyond the meaningful-effect threshold; none is null")
    + HEAD_GAP
  )
  y = (
    add_panel(
      slide,
      RX,
      y,
      [
        head("Distortion, not the budget, predicts outcomes"),
        body(
          "Meaning-preserving sanitization saturates: tightening the budget threefold (ε=3→1) changes distortion and "
          "entity fidelity negligibly (by 0.003 and 0.005). The cliff at ε=6 then shifts every outcome simultaneously."
        ),
        body(
          "Outcomes track distortion, not the nominal budget — so equal-distortion settings are interchangeable, and "
          "looser ones can be pruned before any training compute is spent."
        ),
      ],
    )
    + GAP
  )
  y = add_figure(slide, RX, y, "fig/fig2_cliff.png", 20) + 0.02
  y = add_caption(slide, RX, y, "Equal distortion at ε=1 and ε=3 yields equal outcomes; the cliff is at ε=6") + HEAD_GAP
  y = (
    add_panel(
      slide,
      RX,
      y,
      [
        sub("Safety and privacy, measured separately"),
        body(
          "DP-SGD raises refusal of harmful prompts above the non-private baseline, while text perturbation erodes it "
          "further — safety moves along the same distortion axis as capability, not the privacy budget.",
          bullet=True,
        ),
        body(
          "The failure is fabrication, not corruption: distorted models invent values absent from the record rather "
          "than mis-selecting ones that are present.",
          bullet=True,
        ),
        body(
          "Memorisation stays at the noise floor across every mechanism; the privacy side of the trade rests on the "
          "formal accounting, not on measured leakage.",
          bullet=True,
        ),
      ],
    )
    + HEAD_GAP
  )
  y = (
    add_panel(
      slide,
      RX,
      y,
      [
        head("Mitigation: prune dominated settings"),
        body(
          "Some privacy settings are strictly dominated: another provides stronger privacy at equal distortion and "
          "equal utility. Sanitization (TEM) draws each substitute word w by an exponential mechanism over "
          "embedding-space utility u:"
        ),
      ],
    )
    + 0.14
  )
  y = add_figure(slide, RX, y, "fig/eq_tem.png", 10.5) + 0.16
  y = (
    add_panel(
      slide,
      RX,
      y,
      [
        body(
          "Below a threshold temperature the mechanism saturates — nearly every token is already resampled — so "
          "tightening ε adds no further distortion, and utility, which tracks distortion, likewise remains unchanged."
        ),
        body(
          "Thus ε=3 is dominated by ε=1: identical distortion, equivalent utility (TOST, p < 0.005 on every metric), "
          "yet a 3× stronger guarantee. The full grid can be screened from text alone — retaining the strictest ε on "
          "each plateau — before any fine-tuning compute is committed.",
          bold=True,
        ),
      ],
    )
    + HEAD_GAP
  )
  add_panel(
    slide,
    RX,
    y,
    [
      sub("References"),
      ref("[1] Abadi et al. Deep Learning with Differential Privacy. ACM CCS, 2016."),
      ref("[2] Carvalho et al. TEM: High-Utility Metric Differential Privacy on Text. arXiv, 2021."),
      ref("[3] Ramesh et al. The Privacy–Hallucination Tradeoff in DP Language Models. 2026."),
      ref("[4] Yue et al. Differential Privacy for Text Analytics via Text Sanitization. ACL, 2021."),
      ref("[5] Hu et al. LoRA: Low-Rank Adaptation of Large Language Models. ICLR, 2022."),
      ref("[6] Lin et al. TruthfulQA. ACL, 2022."),
      ref("[7] Zhou et al. Instruction-Following Eval (IFEval). arXiv, 2023."),
      ref("[8] NVIDIA. Nemotron-PII: Synthetic PII/PHI Dataset. 2025."),
    ],
  )


def set_qr(slide: Slide, label_shape: TextFrame, frame_left: int, frame_top: int) -> None:
  for paragraph in list(label_shape.paragraphs):
    paragraph._p.getparent().remove(paragraph._p)
  paragraph = label_shape.add_paragraph()
  run = paragraph.add_run()
  run.text = "Paper & code"
  run.font.size = Pt(16)
  run.font.color.rgb = INK
  left = Emu(Emu(frame_left) + Cm(0.4))
  top = Emu(Emu(frame_top) + Cm(0.4))
  slide.shapes.add_picture("fig/qr.png", left, top, Cm(3.6), Cm(3.6))


def build_poster() -> None:
  presentation = Presentation("template.pptx")
  slide = presentation.slides[0]
  by_id = {shape.shape_id: shape for shape in slide.shapes}
  for shape in list(slide.shapes):
    if shape.shape_id not in {5, 35, 2, 30, 26}:
      shape._element.getparent().remove(shape._element)
  title = by_id[5]
  title.top = Cm(5.2)
  set_title(title.text_frame)
  add_author(slide, 13.5, "Dmitry Scherbakov", "Moscow State University")
  add_author(slide, 30.9, "Fang Liu", "University of Notre Dame")
  divider = by_id[35]
  divider.top = Cm(13.0)
  divider.height = Cm(65.4)
  left_column(slide)
  right_column(slide)
  frame = by_id[2]
  set_qr(slide, by_id[26].text_frame, frame.left, frame.top)
  slide_ids = list(presentation.slides._sldIdLst)
  if len(slide_ids) > 1:
    presentation.part.drop_rel(slide_ids[1].get(qn("r:id")))
    presentation.slides._sldIdLst.remove(slide_ids[1])
  presentation.save("poster.pptx")


if __name__ == "__main__":
  make_figures()
  build_poster()
